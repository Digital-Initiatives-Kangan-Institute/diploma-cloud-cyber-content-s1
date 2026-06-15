#!/usr/bin/env python3
"""Shared YAT brand helpers for building .pptx decks (python-pptx).

Mirrors the .docx brand helpers (build_bc_template.py) for presentations: a 16:9
deck on the YAT palette (brand-pack.md §4) with the simulated-environment
disclosure on every slide's footer (§5.3). Used by the Business Case deck
template and the worked past-engagement deck.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# --- Brand palette (brand-pack.md §4.1) ---
TEAL = "1F5A5C"
TERRACOTTA = "C5613B"
OCHRE = "C99932"
CREAM = "F8F4ED"
CHARCOAL = "1F2329"
GREY = "6B6660"
STONE = "E4DED3"
WHITE = "FFFFFF"
FONT = "Source Sans 3"

DISCLOSURE = "This is a fictional organisation used as a case study in an educational context."

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


def rgb(h):
    return RGBColor.from_string(h)


def new_deck():
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    return prs


def _set_bg(slide, hex_fill=CREAM):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(hex_fill)


def _rect(slide, left, top, width, height, fill_hex):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = rgb(fill_hex)
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _text(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _run(paragraph, text, size, color=CHARCOAL, bold=False, italic=False):
    r = paragraph.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = FONT
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    return r


def _footer(slide):
    # Ochre disclosure bar across the very bottom (brand-pack §5.3).
    bar = _rect(slide, 0, Inches(7.05), EMU_W, Inches(0.45), OCHRE)
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, DISCLOSURE, 9, CHARCOAL, bold=True)


def _title_block(slide, title):
    tb, tf = _text(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(1.0))
    p = tf.paragraphs[0]
    _run(p, title, 28, TEAL, bold=True)
    _rect(slide, Inches(0.72), Inches(1.45), Inches(2.2), Inches(0.05), TERRACOTTA)


def _base(prs, title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide)
    _footer(slide)
    if title:
        _title_block(slide, title)
    return slide


def add_title_slide(prs, title, subtitle, meta_lines):
    slide = _base(prs)
    # Wordmark
    tb, tf = _text(slide, Inches(0.7), Inches(1.6), Inches(11.9), Inches(1.1))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "YAT", 44, TEAL, bold=True)
    _run(p, "  COLLEGE", 24, TERRACOTTA, bold=True)
    # Title
    tb, tf = _text(slide, Inches(0.7), Inches(3.0), Inches(11.9), Inches(1.2))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, title, 40, TEAL, bold=True)
    # Subtitle
    tb, tf = _text(slide, Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.8))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, subtitle, 20, TERRACOTTA, bold=True)
    # Meta
    tb, tf = _text(slide, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.2))
    for i, line in enumerate(meta_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        _run(p, line, 13, GREY)
    return slide


def add_content_slide(prs, title, items, guidance=False):
    """items: list of (text, level) — level 0 or 1. guidance=True greys/italicises."""
    slide = _base(prs, title)
    tb, tf = _text(slide, Inches(0.8), Inches(1.75), Inches(11.7), Inches(5.0))
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(8)
        bullet = "•  " if level == 0 else "–  "
        if guidance:
            _run(p, bullet + text, 15 if level == 0 else 13, GREY, italic=True)
        else:
            _run(p, bullet, 16 if level == 0 else 14, TERRACOTTA, bold=True)
            _run(p, text, 16 if level == 0 else 14, CHARCOAL)
    return slide


def add_statement_slide(prs, title, statement, subtext=""):
    slide = _base(prs, title)
    tb, tf = _text(slide, Inches(1.2), Inches(2.6), Inches(10.9), Inches(2.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    guide = statement.startswith("[")
    _run(p, statement, 30, GREY if guide else TEAL, bold=not guide, italic=guide)
    if subtext:
        tb, tf = _text(slide, Inches(1.6), Inches(5.0), Inches(10.1), Inches(1.4))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, subtext, 15, GREY, italic=subtext.startswith("["))
    return slide


def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = _base(prs, title)
    nrows, ncols = len(rows) + 1, len(headers)
    width = Inches(11.7)
    height = Inches(0.5 * nrows)
    gf = slide.shapes.add_table(nrows, ncols, Inches(0.8), Inches(1.9), width, height)
    table = gf.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    # header
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(TEAL)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        _run(p, h, 12, WHITE, bold=True)
    # body
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(WHITE if r % 2 else CREAM)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            money = val.startswith("$") or val.startswith("+$") or val.startswith("−$")
            if money:
                p.alignment = PP_ALIGN.RIGHT
            _run(p, val, 12, CHARCOAL, bold=(c == 0))
    return slide


def add_closing_slide(prs, line="Questions", subtitle=""):
    slide = _base(prs)
    tb, tf = _text(slide, Inches(0.7), Inches(2.7), Inches(11.9), Inches(1.4))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, line, 40, TEAL, bold=True)
    if subtitle:
        tb, tf = _text(slide, Inches(0.7), Inches(4.1), Inches(11.9), Inches(0.7))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _run(p, subtitle, 16, GREY)
    # small wordmark
    tb, tf = _text(slide, Inches(0.7), Inches(5.0), Inches(11.9), Inches(0.7))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, "YAT", 20, TEAL, bold=True)
    _run(p, "  COLLEGE", 12, TERRACOTTA, bold=True)
    return slide
