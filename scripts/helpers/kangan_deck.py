#!/usr/bin/env python3
"""Kangan-branded slide deck — shared brand + layout helpers.

The single base for every teaching deck. Per-Topic builders do `import kangan_deck as k`
(or `from kangan_deck import *`) and assemble content from these layouts; they hold no brand
or layout code of their own.

Brand reference: documentaion/kangan-branding.md (Kangan/BKI — gold #EDAB0C + charcoal, Roboto;
distinct from the in-world YAT case-study brand).

Layouts provided:
  title_slide      — dark cover (topic number + title + subtitle)
  divider_slide    — section divider (big number on an accent panel)
  content_slide    — title + kicker + bullets (full width)
  visual_slide     — title + kicker + bullets + 0..3 labelled image placeholders
  activity_slide   — gold header band + "ACTIVITY" pill + bullets + timer chip
  demo_slide       — instructor demo (charcoal band + "DEMONSTRATION" pill); sits before an AWS activity
  takeaways_slide  — light-tint "Key takeaways" cards
  table_slide      — title + kicker + branded table (+ optional note)
  close_slide      — dark closing slide (title + lines)
Helpers: new_deck(), save(), placeholder(), and primitives (_rgb/_bg/_rect/_box/_run/_para/
_footer/_bullets).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- brand palette (from kangan-theme CSS tokens) ----
GOLD     = "EDAB0C"   # primary
GOLD_DK  = "D68E10"   # primary-hover / dark
GOLD_BR  = "FBB900"   # bright accent
CHAR     = "2A2929"   # charcoal (secondary)
INK      = "000000"
GREY1    = "484848"
GREY2    = "7A7A7A"
WHITE    = "FFFFFF"
BGLIGHT  = "F9F9F9"
BORDER   = "CCCCCC"
# category accents (per study area / section)
MAGENTA  = "92268F"
SKY      = "27B5CE"
GREEN    = "205F61"
NAVY     = "004488"

FONT_BOLD = "Roboto"
FONT_MED  = "Roboto"
FONT_LT   = "Roboto Light"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# ---------- primitives ----------
def _rgb(h):
    return RGBColor.from_string(h)


def _bg(slide, hexcol):
    f = slide.background.fill
    f.solid(); f.fore_color.rgb = _rgb(hexcol)


def _rect(slide, l, t, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = _rgb(line); sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def _box(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    return tb, tf


def _run(p, text, size, color, bold=False, italic=False, font=FONT_MED):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = _rgb(color)
    return r


def _para(tf, first=False):
    return tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()


def _footer(slide, pageno, accent=GOLD):
    _rect(slide, Inches(0), Inches(7.18), Inches(13.333), Pt(3), fill=accent)
    tb, tf = _box(slide, Inches(0.55), Inches(7.0), Inches(6), Inches(0.35))
    p = _para(tf, True)
    _run(p, "Kangan ", 11, CHAR, bold=True, font=FONT_BOLD)
    _run(p, "Institute", 11, GOLD_DK, bold=True, font=FONT_BOLD)
    tb2, tf2 = _box(slide, Inches(11.5), Inches(7.0), Inches(1.3), Inches(0.35))
    p2 = _para(tf2, True); p2.alignment = PP_ALIGN.RIGHT
    _run(p2, str(pageno), 11, GREY2, font=FONT_MED)


def _bullets(tf, items, base_size=18):
    """items: list of (level, text[, opts]). opts: bold/color/italic/marker/mark_color."""
    for i, it in enumerate(items):
        lvl, text = it[0], it[1]
        opts = it[2] if len(it) > 2 else {}
        p = _para(tf, first=(i == 0))
        p.space_after = Pt(10 if lvl == 0 else 5)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        size = base_size - (3 if lvl >= 1 else 0) - (3 if lvl >= 2 else 0)
        marker = opts.get("marker")
        if marker is None:
            marker = "" if lvl == 0 else ("–  " if lvl == 1 else "·  ")
        indent = Inches(0.0 if lvl == 0 else (0.4 if lvl == 1 else 0.8))
        pPr = p._p.get_or_add_pPr()
        pPr.set('marL', str(int(indent)))
        pPr.set('indent', '0')
        if lvl == 0:
            _run(p, "■  ", size, opts.get("mark_color", GOLD), bold=True, font=FONT_BOLD)
        elif marker:
            _run(p, marker, size, GREY2, font=FONT_MED)
        _run(p, text, size, opts.get("color", CHAR),
             bold=opts.get("bold", False), italic=opts.get("italic", False),
             font=FONT_BOLD if opts.get("bold") else FONT_MED)


# ---------- deck + save ----------
def new_deck():
    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H
    return prs


def save(prs, path):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    print(f"Wrote {path} ({len(prs.slides._sldIdLst)} slides)")


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- layouts ----------
def title_slide(prs, topic_no, title, subtitle):
    s = _blank(prs)
    _bg(s, CHAR)
    _rect(s, Inches(0), Inches(0), Inches(0.45), Inches(7.5), fill=GOLD)
    tb, tf = _box(s, Inches(1.0), Inches(2.2), Inches(11.5), Inches(2.6))
    p = _para(tf, True)
    _run(p, f"TOPIC {topic_no}", 20, GOLD_BR, bold=True, font=FONT_BOLD)
    p2 = tf.add_paragraph(); p2.space_before = Pt(8)
    _run(p2, title, 46, WHITE, bold=True, font=FONT_BOLD)
    p3 = tf.add_paragraph(); p3.space_before = Pt(12)
    _run(p3, subtitle, 22, "D1D1D1", italic=True, font=FONT_LT)
    tb2, tf2 = _box(s, Inches(1.0), Inches(6.5), Inches(8), Inches(0.5))
    pp = _para(tf2, True)
    _run(pp, "Kangan ", 16, WHITE, bold=True, font=FONT_BOLD)
    _run(pp, "Institute", 16, GOLD, bold=True, font=FONT_BOLD)
    return s


def divider_slide(prs, number, title, kicker, accent):
    s = _blank(prs)
    _bg(s, WHITE)
    _rect(s, Inches(0), Inches(0), Inches(4.3), Inches(7.5), fill=accent)
    tb, tf = _box(s, Inches(0.4), Inches(2.3), Inches(3.6), Inches(2.4))
    p = _para(tf, True)
    _run(p, number, 130, WHITE, bold=True, font=FONT_BOLD)
    tb2, tf2 = _box(s, Inches(4.9), Inches(2.7), Inches(7.8), Inches(2.4))
    p1 = _para(tf2, True)
    _run(p1, "SECTION", 16, accent, bold=True, font=FONT_BOLD)
    p2 = tf2.add_paragraph(); p2.space_before = Pt(6)
    _run(p2, title, 38, CHAR, bold=True, font=FONT_BOLD)
    if kicker:
        p3 = tf2.add_paragraph(); p3.space_before = Pt(10)
        _run(p3, kicker, 20, GREY1, italic=True, font=FONT_LT)
    _footer(s, "", accent=accent)
    return s


def _title_block(prs, title, kicker, accent):
    s = _blank(prs)
    _bg(s, WHITE)
    tb, tf = _box(s, Inches(0.7), Inches(0.5), Inches(11.9), Inches(1.05))
    p = _para(tf, True)
    _run(p, title, 28, CHAR, bold=True, font=FONT_BOLD)
    if kicker:
        pk = tf.add_paragraph(); pk.space_before = Pt(2)
        _run(pk, kicker, 18, GOLD_DK, italic=True, font=FONT_LT)
    _rect(s, Inches(0.72), Inches(1.5), Inches(1.6), Pt(4), fill=accent)
    return s


def content_slide(prs, pageno, title, kicker, bullets, accent=GOLD, base=18):
    s = _title_block(prs, title, kicker, accent)
    tb2, tf2 = _box(s, Inches(0.72), Inches(1.85), Inches(11.9), Inches(4.9))
    _bullets(tf2, bullets, base_size=base)
    _footer(s, pageno, accent=accent)
    return s


def placeholder(slide, l, t, w, h, label, accent=GOLD_DK):
    box = _rect(slide, l, t, w, h, fill="FBF6E6", line=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    ln = box.line._get_or_add_ln()
    ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    tb, tf = _box(slide, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _run(p, "🖼", 30, accent, font=FONT_BOLD)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(4)
    _run(p2, "IMAGE", 12, accent, bold=True, font=FONT_BOLD)
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(2)
    _run(p3, label, 12, GREY1, italic=True, font=FONT_LT)


def visual_slide(prs, pageno, title, kicker, bullets, images, accent=GOLD):
    """Content slide with 0..3 labelled image placeholders + bullets."""
    s = _title_block(prs, title, kicker, accent)
    n = len(images)
    if n == 0:
        tb, tf = _box(s, Inches(0.72), Inches(1.85), Inches(11.9), Inches(4.9))
        _bullets(tf, bullets, base_size=18)
    elif n == 1:
        if bullets:
            tb, tf = _box(s, Inches(0.72), Inches(1.9), Inches(6.0), Inches(4.7))
            _bullets(tf, bullets, base_size=17)
            placeholder(s, Inches(7.0), Inches(1.95), Inches(5.6), Inches(4.4), images[0], accent)
        else:
            placeholder(s, Inches(2.7), Inches(1.95), Inches(7.9), Inches(4.5), images[0], accent)
    elif n == 2:
        if bullets:
            tb, tf = _box(s, Inches(0.72), Inches(1.85), Inches(11.9), Inches(1.3))
            _bullets(tf, bullets, base_size=16)
            ytop = Inches(3.35)
        else:
            ytop = Inches(2.1)
        placeholder(s, Inches(0.9), ytop, Inches(5.5), Inches(3.3), images[0], accent)
        placeholder(s, Inches(6.9), ytop, Inches(5.5), Inches(3.3), images[1], accent)
    else:  # 3 across
        if bullets:
            tb, tf = _box(s, Inches(0.72), Inches(1.85), Inches(11.9), Inches(1.0))
            _bullets(tf, bullets, base_size=16)
            ytop = Inches(3.0)
        else:
            ytop = Inches(2.2)
        for x, lab in zip([Inches(0.72), Inches(4.42), Inches(8.12)], images):
            placeholder(s, x, ytop, Inches(3.6), Inches(3.4), lab, accent)
    _footer(s, pageno, accent=accent)
    return s


def activity_slide(prs, pageno, title, bullets, timer, accent=GOLD):
    s = _blank(prs)
    _bg(s, WHITE)
    _rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.35), fill=accent)
    pill = _rect(s, Inches(0.7), Inches(0.42), Inches(1.85), Inches(0.5),
                 fill=CHAR, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pa = pill.text_frame.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
    _run(pa, "ACTIVITY", 14, GOLD_BR, bold=True, font=FONT_BOLD)
    tb, tf = _box(s, Inches(2.8), Inches(0.34), Inches(9.8), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    _run(_para(tf, True), title, 26, WHITE, bold=True, font=FONT_BOLD)
    tb2, tf2 = _box(s, Inches(0.72), Inches(1.7), Inches(11.9), Inches(4.7))
    _bullets(tf2, bullets, base_size=17)
    _rect(s, Inches(0.7), Inches(6.45), Inches(4.4), Inches(0.5),
          fill=BGLIGHT, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tbt, tft = _box(s, Inches(0.9), Inches(6.5), Inches(4.0), Inches(0.4), anchor=MSO_ANCHOR.MIDDLE)
    _run(_para(tft, True), "⏱  " + timer, 15, CHAR, bold=True, font=FONT_BOLD)
    _footer(s, pageno, accent=accent)
    return s


def demo_slide(prs, pageno, title, bullets, accent=GOLD, source=None):
    """Demo of an AWS task — sits right before the matching activity.
    The AWS-practical flow is teach → DEMONSTRATE → practice (activity). Charcoal header
    band (vs the activity's accent band) so 'watch' vs 'do' read differently at a glance.
    `source` = an AWS recorded-demo reference (e.g. "ACF M04 · IAM"): when given, the slide
    cues the recorded demo (RECORDED DEMO pill); when None, it's a live instructor demo."""
    recorded = bool(source)
    s = _blank(prs)
    _bg(s, WHITE)
    _rect(s, Inches(0), Inches(0), Inches(13.333), Inches(1.35), fill=CHAR)
    pill_w = Inches(2.55) if recorded else Inches(2.45)
    pill = _rect(s, Inches(0.7), Inches(0.42), pill_w, Inches(0.5),
                 fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    pa = pill.text_frame.paragraphs[0]; pa.alignment = PP_ALIGN.CENTER
    _run(pa, "RECORDED DEMO" if recorded else "DEMONSTRATION", 13, CHAR, bold=True, font=FONT_BOLD)
    tb, tf = _box(s, Inches(3.5), Inches(0.34), Inches(9.1), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    _run(_para(tf, True), title, 25, WHITE, bold=True, font=FONT_BOLD)
    tb2, tf2 = _box(s, Inches(0.72), Inches(1.7), Inches(11.9), Inches(4.7))
    _bullets(tf2, bullets, base_size=17)
    if recorded:
        chip = "▶  Play the AWS recorded demo (%s) — you do this next" % source
    else:
        chip = "▶  Instructor demonstrates live (no recorded demo) — you do this next"
    chip_w = Inches(8.6) if recorded else Inches(7.2)
    _rect(s, Inches(0.7), Inches(6.45), chip_w, Inches(0.5),
          fill=BGLIGHT, line=BORDER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tbt, tft = _box(s, Inches(0.9), Inches(6.5), Emu(int(chip_w) - int(Inches(0.4))), Inches(0.4),
                    anchor=MSO_ANCHOR.MIDDLE)
    _run(_para(tft, True), chip, 14, CHAR, bold=True, font=FONT_BOLD)
    _footer(s, pageno, accent=accent)
    return s


def takeaways_slide(prs, pageno, section_label, points, accent=GOLD):
    s = _blank(prs)
    _bg(s, BGLIGHT)
    tb, tf = _box(s, Inches(0.7), Inches(0.6), Inches(11.9), Inches(1.2))
    p = _para(tf, True)
    _run(p, "Key takeaways", 32, CHAR, bold=True, font=FONT_BOLD)
    pk = tf.add_paragraph(); pk.space_before = Pt(2)
    _run(pk, section_label, 18, accent, bold=True, font=FONT_BOLD)
    _rect(s, Inches(0.72), Inches(1.75), Inches(1.6), Pt(4), fill=accent)
    top, gap, card_h = Inches(2.15), Inches(0.12), Inches(0.92)
    for i, txt in enumerate(points):
        y = Emu(int(top) + i * (int(card_h) + int(gap)))
        _rect(s, Inches(0.7), y, Inches(11.95), card_h, fill=WHITE, line=BORDER)
        _rect(s, Inches(0.7), y, Inches(0.12), card_h, fill=accent)
        tbc, tfc = _box(s, Inches(1.05), y, Inches(11.4), card_h, anchor=MSO_ANCHOR.MIDDLE)
        pc = _para(tfc, True); pc.line_spacing = 1.0
        _run(pc, txt, 17, CHAR, font=FONT_MED)
    _footer(s, pageno, accent=accent)
    return s


def table_slide(prs, pageno, title, kicker, headers, rows, accent=GOLD, col_widths=None, note=None):
    s = _title_block(prs, title, kicker, accent)
    nrows, ncols = len(rows) + 1, len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.72), Inches(2.0),
                            Inches(11.9), Inches(0.5 * nrows)).table
    gt.first_row = False; gt.horz_banding = False
    if col_widths:
        for c, w in enumerate(col_widths):
            gt.columns[c].width = Inches(w)
    for c, htext in enumerate(headers):
        cell = gt.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(accent)
        cell.margin_left = cell.margin_right = Pt(6)
        _run(cell.text_frame.paragraphs[0], htext, 14, WHITE, bold=True, font=FONT_BOLD)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(WHITE if r % 2 else BGLIGHT)
            cell.margin_left = cell.margin_right = Pt(6)
            cell.margin_top = cell.margin_bottom = Pt(4)
            bold = (c == 0 and ncols == 2 and headers[0] == "Field")
            _run(cell.text_frame.paragraphs[0], val, 13, CHAR, bold=bold,
                 font=FONT_BOLD if bold else FONT_MED)
    if note:
        tbn, tfn = _box(s, Inches(0.72), Inches(6.5), Inches(11.9), Inches(0.5))
        _run(_para(tfn, True), note, 15, GREY1, italic=True, font=FONT_LT)
    _footer(s, pageno, accent=accent)
    return s


def close_slide(prs, title, lines, accent=GOLD):
    s = _blank(prs)
    _bg(s, CHAR)
    _rect(s, Inches(0), Inches(0), Inches(0.45), Inches(7.5), fill=accent)
    tb, tf = _box(s, Inches(1.0), Inches(1.6), Inches(11.5), Inches(4.4))
    _run(_para(tf, True), title, 40, WHITE, bold=True, font=FONT_BOLD)
    for txt in lines:
        pp = tf.add_paragraph(); pp.space_before = Pt(12); pp.line_spacing = 1.05
        _run(pp, txt, 19, "E8E8E8", font=FONT_LT)
    return s
